"""Build the SIH 2026 idea deck: template chrome, our content, our diagrams.

Two passes, because neither library does both jobs well.

1. **Text, by editing the packed XML.** `python-pptx` collapses a paragraph to a
   single unstyled run on assignment, which loses the template's fonts and
   bullets. Each content box is rebuilt *whole* from `SPEC` -- every paragraph
   regenerated, never patched -- so the script is idempotent and can read its own
   output as input. That also fixes a defect in the earlier version, which
   replaced only the first `<a:t>` of a prompt and so left the tail of any
   multi-run prompt stranded at the end of our bullet ("... not as a win. for
   overcoming these challenges").

2. **Layout and graphics, through python-pptx.** Diagrams are native shapes, not
   pictures, so they stay sharp and remain editable. Every shape this script adds
   is named `AS_*` and deleted before it is re-added, so pass 2 is idempotent too.

The template's idea pointers may not be changed, so all of them are kept verbatim
as bold sub-headings, with our content beneath the prompt it answers.

    python scripts/build_sih_deck.py

Needs `python-pptx` and `Pillow`; both live outside the project venv, so run it
with the system interpreter.
"""

from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "docs" / "Airshed-SIH2026-Idea.pptx"
WORK = ROOT / "data" / "_deck_build"          # gitignored scratch
SHOTS = ROOT / "docs" / "results"

# Screenshots of the running dashboard, cropped to the cards we want on slide 4.
RAW_LIVE = SHOTS / "screenshot_forecast_2026-08-25.png"
RAW_MAP = SHOTS / "screenshot_surface_2026-08-25.png"
CROP_LIVE = SHOTS / "dashboard_live.png"
CROP_MAP = SHOTS / "dashboard_surface.png"

HEADER_PT, PROMPT_PT, BODY_PT = 1800, 1350, 1150

# Slide 2's title placeholder ships as "IDEA TITLE" and is ours to fill. Keep it
# short: the placeholder is a wide serif and wraps into the content box.
IDEA_TITLE = "AIRSHED — 72-HOUR GRAP FORECAST"
IDEA_TITLE_PT = 2800

# ---------------------------------------------------------------- content ---
# ("h", ...) underlined section pointer   ("p", ...) template prompt, kept verbatim
# ("b", ...) our content, one point per line

TITLE_LINES = [
    "Problem Statement ID – SIH26082",
    "Problem Statement Title- Air Pollution–Weather Coupled Forecasting "
    "System for Delhi NCR",
    "Theme- Disaster Management",
    "PS Category- Software",
    "Team ID- <your team ID>",
    "Team Name (Registered on portal)- <your team name>",
]

SPEC: dict[int, list[tuple[str, str]]] = {
    2: [
        ("h", "Proposed Solution (Describe your Idea/Solution/Prototype)"),
        ("p", "Detailed explanation of the proposed solution"),
        ("b", "Correct a physics forecast instead of rebuilding one. CAMS "
              "(Copernicus) publishes global PM2.5 seven days ahead, free and "
              "keyless — but at ~40 km, and biased over India."),
        ("b", "CAMS supplies the future: winds, transport, regional build-up. "
              "Our layer supplies the local: station-level bias correction and "
              "sub-grid detail."),
        ("b", "Trained on 1.07 M station-hours, 77 CPCB stations, Feb 2025 "
              "onward. Every input is an archived forecast, never reanalysis, so "
              "training and serving see the same distribution."),
        ("p", " How it addresses the problem"),
        ("b", "72 h PM2.5 at every station with a 10–90% interval, jointly with "
              "boundary-layer height and visibility — pollution and weather as "
              "one coupled system, which is what the problem statement asks for."),
        ("b", "The output is the decision, not a number: the probability each "
              "GRAP stage I–IV is reached, with lead time and driver "
              "attribution (\"upwind fires + shallow inversion\")."),
        ("b", "Stage III is caught at 73.8% recall with 72 h of lead. GRAP "
              "measures need days to take effect — that is the difference "
              "between a forecast and a warning."),
        ("p", "Innovation and uniqueness of the solution "),
        ("b", "The baseline is external and published, so beating it is a "
              "measured result, not a claim: 31.5% below raw CAMS and 20.6% "
              "below persistence, on 5 of 5 seasonal folds."),
        ("b", "Persistence — \"tomorrow is like today\" — appears in every "
              "results table. It beats most published AQ models, and most "
              "published AQ models never report it."),
        ("b", "Null results are published beside the wins: fires +0.9%, upwind "
              "corridor +0.0%, coupled multi-output −0.1%, each inside its own "
              "scatter. The same correction bolts onto the ministry's own model."),
    ],
    3: [
        ("p", "Technologies to be used (e.g. programming languages, frameworks, "
              "hardware)"),
        ("b", "Python 3.11 · LightGBM · Polars · DuckDB over Parquet · FastAPI · "
              "MapLibre GL + deck.gl · APScheduler — no GPU, no database server, "
              "no paid API or key for the physics: it runs on one laptop."),
        ("p", "Methodology and process for implementation (Flow Charts/Images/ "
              "working prototype)"),
        ("b", "Below, left to right. Direct multi-horizon quantile heads (no "
              "recursive rollout, so 72 h error cannot compound), conformal "
              "calibration at 85.1% coverage, and a cache the demo never reads "
              "past — an upstream outage degrades it instead of killing it."),
    ],
    4: [
        ("p", "Analysis of the feasibility of the idea"),
        ("b", "Built and running end-to-end, not a proposal: 9.1 M cached rows, "
              "77 stations, live 72 h forecast, historical replay, GRAP "
              "probabilities and a downscaled surface — the panels at right are "
              "that system, running today."),
        ("b", "Zero data cost and zero licence cost. A resident archive job "
              "re-checks every 30 minutes and asks whether today's run is in the "
              "store, so sleep and reboots cannot make it miss."),
        ("p", "Potential challenges and risks"),
        ("b", "One winter of trainable ground truth, so effects of a few percent "
              "cannot yet be separated from seasonal noise."),
        ("b", "Train/serve gap: archived CAMS and live CAMS differ by −32 µg/m³ "
              "at day 2 on 3 run days; 20 settled days are needed before a "
              "correction for it can honestly be fitted."),
        ("b", "CPCB stations go offline and relocate, and a nine-day January 2026 "
              "hole is upstream and unrecoverable."),
        ("p", "Strategies for overcoming these challenges"),
        ("b", "Persistence in every table, whole-season holdout with a 96 h "
              "embargo, and 5 rolling-origin folds — a gain smaller than its own "
              "spread is reported as \"cannot tell\", not as a win."),
        ("b", "Station gaps are treated as missing, never interpolated, and the "
              "CAMS run gap is shown on the dashboard as a \"known bias\" note "
              "instead of being quietly absorbed into the forecast."),
        ("b", "Everything is cached and re-runnable offline; November 2026 is a "
              "second episode season that settles the open questions in one run."),
    ],
    5: [
        ("p", "Potential impact on the target audience"),
        ("b", "CAQM and DPCC act through GRAP, and GRAP takes days to bite. "
              "Stage III at 73.8% recall and 66.0% precision with 72 h of lead "
              "turns a forecast into an actionable warning."),
        ("b", "30 M+ NCR residents get a neighbourhood answer, not a city "
              "average: a 0.05° surface, 81.8 µg/m³ leave-one-station-out, "
              "24.6% better than inverse-distance interpolation."),
        ("b", "Every number ships with its interval and a visible \"last synced\" "
              "age, so an operator can see when the system is unsure or stale — "
              "the failure mode that ends operational trust in a forecast."),
        ("p", "Benefits of the solution (social, economic, environmental, etc.)"),
        ("b", "Social — earlier school closures, construction halts and health "
              "advisories, with the interval shown everywhere so a marginal call "
              "looks marginal instead of certain."),
        ("b", "Economic — fog risk three days out for IGI diversions and highway "
              "pile-ups: visibility skill improves 7.5% from pollution inputs "
              "(we report the number; the mechanism is not yet established)."),
        ("b", "Environmental — the surface says which part of the city is worse, "
              "and every cell carries its distance to the nearest monitor, so an "
              "unmonitored block is marked uncertain rather than quietly guessed."),
        ("b", "Institutional — the correction layer is model-agnostic. It applies "
              "to the ministry's own SAFAR / Delhi DSS output as easily as to "
              "CAMS, so it adds to existing investment rather than competing "
              "with it, and it costs nothing to run."),
    ],
    6: [
        ("p", "Details / Links of the reference and research work"),
        ("b", "CAMS global forecast (Copernicus Atmosphere Monitoring Service), "
              "read free and keyless through the Open-Meteo Air Quality API — "
              "air-quality-api.open-meteo.com."),
        ("b", "Open-Meteo Historical Forecast and Previous Runs APIs for archived "
              "past forecast runs — the training meteorology, at the lead it is "
              "actually used at. ERA5 reanalysis is deliberately excluded: it "
              "lags real time by ~5 days and cannot be served."),
        ("b", "CPCB CAAQMS ground truth via the OpenAQ v3 API (api.openaq.org); "
              "NASA FIRMS VIIRS/MODIS active fire detections for stubble burning "
              "(firms.modaps.eosdis.nasa.gov); Iowa State IEM ASOS archive for "
              "measured METAR visibility at VIDP (mesonet.agron.iastate.edu)."),
        ("b", "GRAP stage thresholds transcribed from the CAQM published "
              "schedule; CPCB National AQI breakpoints for the 24-hour mean."),
        ("b", "Prior art we position against, not around: SAFAR (IITM / MoES) and "
              "the Delhi Decision Support System already run the physics for "
              "NCR. Correcting an operational forecast is standard practice in "
              "operational meteorology; applying it to CAMS over NCR, and "
              "outputting the GRAP decision, is ours."),
        ("b", "Method and every number on these slides: docs/results/ in the "
              "project repo — ablation.md, rolling.md, grap.md, loso.md, "
              "coupling.md, leadmatch.md, each regenerable by one command."),
    ],
}

# ------------------------------------------------------------- pass 1: XML ---

ARIAL = ('<a:latin typeface="Arial" pitchFamily="34" charset="0"/>'
         '<a:cs typeface="Arial" pitchFamily="34" charset="0"/>')
WINGDINGS = ('<a:buFont typeface="Wingdings" panose="05000000000000000000" '
             'pitchFamily="2" charset="2"/><a:buChar char="v"/>')
BULLET = ('<a:buFont typeface="Arial" panose="020B0604020202020204" '
          'pitchFamily="34" charset="0"/><a:buChar char="•"/>')
DASH = ('<a:buFont typeface="Arial" panose="020B0604020202020204" '
        'pitchFamily="34" charset="0"/><a:buChar char="–"/>')


def esc(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def para(kind: str, text: str) -> str:
    """One paragraph of the content box, built from scratch in the template's style."""
    if kind == "h":
        ppr = f'<a:pPr marL="342900" indent="-342900">{WINGDINGS}</a:pPr>'
        rpr = (f'<a:rPr lang="en-US" sz="{HEADER_PT}" b="1" u="sng" dirty="0">'
               f'<a:solidFill><a:schemeClr val="tx2"/></a:solidFill>{ARIAL}</a:rPr>')
    elif kind == "p":
        ppr = f'<a:pPr marL="285750" indent="-285750" algn="just">{BULLET}</a:pPr>'
        rpr = f'<a:rPr lang="en-US" sz="{PROMPT_PT}" b="1" dirty="0">{ARIAL}</a:rPr>'
    else:
        ppr = (f'<a:pPr marL="514350" indent="-228600" algn="just">'
               f'<a:lnSpc><a:spcPct val="97000"/></a:lnSpc>'
               f'<a:spcBef><a:spcPts val="300"/></a:spcBef>{DASH}</a:pPr>')
        rpr = f'<a:rPr lang="en-US" sz="{BODY_PT}" dirty="0">{ARIAL}</a:rPr>'
    return f"<a:p>{ppr}<a:r>{rpr}<a:t>{esc(text)}</a:t></a:r></a:p>"


def set_title(xml: str, text: str, sz: int) -> str:
    """Rewrite the slide's title placeholder as a single sized run.

    The placeholder ships as an empty run, a line break and then "IDEA TITLE",
    and it does not wrap: at the template's 36 pt anything much longer runs under
    the team-name oval on one side and the SIH mark on the other. Replacing the
    whole paragraph sets the size and drops the blank first line together.
    """
    m = re.search(r'name="Title 1".*?</p:txBody>', xml, re.S)
    if not m:
        raise SystemExit("  ! no Title 1 placeholder")
    box = m.group(0)
    times = ('<a:latin typeface="Times New Roman" '
             'panose="02020603050405020304" pitchFamily="18" charset="0"/>'
             '<a:cs typeface="Times New Roman" '
             'panose="02020603050405020304" pitchFamily="18" charset="0"/>')
    para_xml = (f'<a:p><a:pPr eaLnBrk="1" hangingPunct="1"/><a:r>'
                f'<a:rPr lang="en-US" sz="{sz}" b="1" dirty="0">{times}</a:rPr>'
                f"<a:t>{esc(text)}</a:t></a:r></a:p>")
    head = box[: box.index("<a:p>")]
    return xml.replace(box, head + para_xml + "</p:txBody>", 1)


def rebuild_box(xml: str, shape_name: str, paras: list[str]) -> str:
    m = re.search(rf'name="{shape_name}".*?</p:txBody>', xml, re.S)
    if not m:
        raise SystemExit(f"  ! shape not found: {shape_name}")
    box = m.group(0)
    head = box[: box.index("<a:p>")]
    return xml.replace(box, head + "".join(paras) + "</p:txBody>", 1)


def fill_text() -> Path:
    if WORK.exists():
        shutil.rmtree(WORK)
    with zipfile.ZipFile(DECK) as z:
        z.extractall(WORK)

    p1 = WORK / "ppt" / "slides" / "slide1.xml"
    title = [
        f'<a:p><a:pPr algn="l"/><a:r><a:rPr lang="en-US" sz="1600" dirty="0">'
        f"{ARIAL}</a:rPr><a:t>{esc(line)}</a:t></a:r></a:p>"
        for line in TITLE_LINES
    ]
    p1.write_text(rebuild_box(p1.read_text(encoding="utf-8"), "TextBox 9", title),
                  encoding="utf-8")
    print("  slide 1: title page")

    for n, items in SPEC.items():
        p = WORK / "ppt" / "slides" / f"slide{n}.xml"
        paras = [para(kind, text) for kind, text in items]
        xml = rebuild_box(p.read_text(encoding="utf-8"), "TextBox 8", paras)
        if n == 2:
            xml = set_title(xml, IDEA_TITLE, IDEA_TITLE_PT)
        p.write_text(xml, encoding="utf-8")
        kept = sum(1 for k, _ in items if k in "hp")
        print(f"  slide {n}: {kept} template pointers kept, "
              f"{len(items) - kept} points")

    out = WORK.parent / "_deck_stage1.pptx"
    out.unlink(missing_ok=True)
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for f in sorted(WORK.rglob("*")):
            if f.is_file():
                z.write(f, f.relative_to(WORK).as_posix())
    return out


# ------------------------------------------------------- pass 2: geometry ---

INK = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2F, 0x6F, 0xB2)
PALE = RGBColor(0xEE, 0xF3, 0xF9)
EDGE = RGBColor(0xA9, 0xBC, 0xD0)
TEAL = RGBColor(0x1E, 0x7A, 0x66)
TEAL_BG = RGBColor(0xE3, 0xF1, 0xEC)
AMBER = RGBColor(0xB2, 0x6A, 0x00)
AMBER_BG = RGBColor(0xFD, 0xF1, 0xDF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x6B, 0x7D)


def clear_ours(slide) -> None:
    for shp in list(slide.shapes):
        if shp.name.startswith("AS_"):
            shp._element.getparent().remove(shp._element)


def move(slide, name: str, x, y, w, h) -> None:
    for shp in slide.shapes:
        if shp.name == name:
            shp.left, shp.top, shp.width, shp.height = (
                Inches(x), Inches(y), Inches(w), Inches(h))
            return
    raise SystemExit(f"  ! cannot move {name}")


def textbox(slide, x, y, w, h, lines, *, name, align=PP_ALIGN.CENTER,
            color=INK, anchor=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    for i, (text, pt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.name = "Arial"
        r.font.color.rgb = col or color
    return box


def node(slide, x, y, w, h, lines, *, name, fill, edge, radius=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.name = name
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = edge
    s.line.width = Pt(1.0)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, (text, pt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 0.95
        r = p.add_run()
        r.text = text
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.name = "Arial"
        r.font.color.rgb = col
    return s


def arrow(slide, x, y, w, h, *, name, color=BLUE, down=False):
    s = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(y), Inches(w), Inches(h))
    s.name = name
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def crop(src: Path, dst: Path, box: tuple[int, int, int, int]) -> Path:
    if not src.exists():
        raise SystemExit(f"  ! screenshot missing: {src}")
    Image.open(src).convert("RGB").crop(box).save(dst)
    return dst


# ------------------------------------------------------------- the slides ---

def slide2(s) -> None:
    """Left column of text; right column a four-stage schematic of the stack."""
    clear_ours(s)
    move(s, "TextBox 8", 0.50, 1.28, 7.55, 5.30)
    x, w = 8.32, 4.62
    node(s, x, 1.42, w, 0.42,
         [("THE CORRECTION STACK", 11, True, WHITE)],
         name="AS_hdr", fill=INK, edge=INK, radius=0.16)
    node(s, x, 2.00, w, 0.86,
         [("CAMS global physics forecast", 12, True, INK),
          ("~40 km · free · 7 days ahead · biased over India", 9.5, False, MUTED)],
         name="AS_n1", fill=PALE, edge=EDGE)
    arrow(s, x + w / 2 - 0.14, 2.92, 0.28, 0.26, name="AS_a1", down=True)
    node(s, x, 3.24, w, 0.86,
         [("Local signal we add", 12, True, INK),
          ("77 CPCB stations · GFS wind & boundary layer · FIRMS fires · "
           "METAR visibility", 9.5, False, MUTED)],
         name="AS_n2", fill=PALE, edge=EDGE)
    arrow(s, x + w / 2 - 0.14, 4.16, 0.28, 0.26, name="AS_a2", down=True)
    node(s, x, 4.48, w, 0.86,
         [("Airshed correction layer", 12, True, WHITE),
          ("LightGBM quantile · direct 24 / 48 / 72 h heads", 9.5, False,
           RGBColor(0xC9, 0xDA, 0xEC))],
         name="AS_n3", fill=INK, edge=INK)
    arrow(s, x + w / 2 - 0.14, 5.40, 0.28, 0.26, name="AS_a3", color=TEAL, down=True)
    node(s, x, 5.72, w, 0.86,
         [("PM2.5 p10–p50–p90  ·  0.05° surface", 11.5, True, TEAL),
          ("GRAP stage probability + driver attribution", 9.5, False, TEAL)],
         name="AS_n4", fill=TEAL_BG, edge=TEAL)


def slide3(s) -> None:
    """Text band across the top; the pipeline flow chart underneath it."""
    clear_ours(s)
    move(s, "TextBox 8", 0.50, 1.26, 12.35, 1.55)

    top, bot = 3.14, 6.58
    height = bot - top
    left, right = 0.44, 12.90
    gap = 0.40
    cw = (right - left - 4 * gap) / 5
    xs = [left + i * (cw + gap) for i in range(5)]
    mid = top + height / 2

    def stage(i, text):
        textbox(s, xs[i], top - 0.28, cw, 0.24,
                [(text, 9, True, BLUE)], name=f"AS_stage{i}")

    # 1 — sources
    stage(0, "SOURCES  (all free)")
    srcs = [
        ("CAMS PM2.5 forecast", "Open-Meteo · ~40 km · 7 d"),
        ("GFS meteorology", "wind, boundary layer, humidity"),
        ("CPCB CAAQMS ground truth", "via OpenAQ · 77 NCR stations"),
        ("NASA FIRMS fires", "VIIRS/MODIS · upwind load"),
        ("METAR visibility", "IEM archive · VIDP hourly"),
    ]
    bh = (height - 4 * 0.09) / 5
    for i, (a, b) in enumerate(srcs):
        node(s, xs[0], top + i * (bh + 0.09), cw, bh,
             [(a, 9.5, True, INK), (b, 8, False, MUTED)],
             name=f"AS_src{i}", fill=PALE, edge=EDGE, radius=0.14)
    arrow(s, xs[0] + cw + 0.06, mid - 0.11, gap - 0.12, 0.22, name="AS_ar0")

    # 2 — ingest and cache
    stage(1, "INGEST → CACHE")
    node(s, xs[1], top + 0.55, cw, height - 1.10,
         [("One module per source", 10, True, INK),
          ("fetch(start, end) → Parquet,", 8.5, False, MUTED),
          ("partitioned by date, idempotent", 8.5, False, MUTED),
          ("", 5, False, MUTED),
          ("DuckDB over Parquet · 9.1 M rows", 8.5, False, MUTED),
          ("Hourly UTC index · IST only in the UI", 8.5, False, MUTED),
          ("", 5, False, MUTED),
          ("Nothing below this line calls the network", 8.5, True, BLUE)],
         name="AS_ingest", fill=WHITE, edge=BLUE, radius=0.08)
    arrow(s, xs[1] + cw + 0.06, mid - 0.11, gap - 0.12, 0.22, name="AS_ar1")

    # 3 — features
    stage(2, "FEATURES")
    node(s, xs[2], top + 0.55, cw, height - 1.10,
         [("Aligned supervised table", 10, True, INK),
          ("1.07 M station-hours", 8.5, False, MUTED),
          ("", 5, False, MUTED),
          ("CAMS at the target hour", 8.5, False, MUTED),
          ("Observation lags & rolling means", 8.5, False, MUTED),
          ("Meteorology at true forecast lead", 8.5, False, MUTED),
          ("Upwind fire load · calendar & season", 8.5, False, MUTED),
          ("", 5, False, MUTED),
          ("Split by time block, never at random", 8.5, True, BLUE)],
         name="AS_feat", fill=WHITE, edge=BLUE, radius=0.08)
    arrow(s, xs[2] + cw + 0.06, mid - 0.11, gap - 0.12, 0.22, name="AS_ar2")

    # 4 — model and the baselines it must beat
    stage(3, "CORRECT")
    node(s, xs[3], top, cw, 2.05,
         [("LightGBM corrector", 10.5, True, WHITE),
          ("9 quantile heads —", 8.5, False, RGBColor(0xC9, 0xDA, 0xEC)),
          ("24/48/72 h × p10/p50/p90", 8.5, False, RGBColor(0xC9, 0xDA, 0xEC)),
          ("Direct heads: no recursive rollout", 8.5, False,
           RGBColor(0xC9, 0xDA, 0xEC)),
          ("Conformal calibration → 85.1% coverage", 8.5, False,
           RGBColor(0xC9, 0xDA, 0xEC))],
         name="AS_model", fill=INK, edge=INK, radius=0.08)
    node(s, xs[3], top + 2.20, cw, height - 2.20,
         [("Baselines & evaluation", 9.5, True, AMBER),
          ("persistence · raw CAMS · scaled CAMS", 8, False, AMBER),
          ("5 rolling-origin folds, whole-season holdout", 8, False, AMBER)],
         name="AS_eval", fill=AMBER_BG, edge=AMBER, radius=0.10)
    arrow(s, xs[3] + cw + 0.06, mid - 0.11, gap - 0.12, 0.22, name="AS_ar3",
          color=TEAL)

    # 5 — decision and serving
    stage(4, "DECIDE → SERVE")
    oh = (height - 2 * 0.13) / 3
    node(s, xs[4], top, cw, oh,
         [("GRAP stage probability", 9.5, True, TEAL),
          ("P(stage I–IV) + lead time; SHAP grouped", 8, False, TEAL),
          ("into human causes", 8, False, TEAL)],
         name="AS_out0", fill=TEAL_BG, edge=TEAL, radius=0.10)
    node(s, xs[4], top + oh + 0.13, cw, oh,
         [("Downscaled surface", 9.5, True, TEAL),
          ("Wind-aware graph → 0.05° grid, each cell", 8, False, TEAL),
          ("carrying its distance to a monitor", 8, False, TEAL)],
         name="AS_out1", fill=TEAL_BG, edge=TEAL, radius=0.10)
    node(s, xs[4], top + 2 * (oh + 0.13), cw, oh,
         [("FastAPI + MapLibre dashboard", 9.5, True, TEAL),
          ("/forecast /replay /grap · scheduled refresh", 8, False, TEAL),
          ("with a visible \"last synced\" age", 8, False, TEAL)],
         name="AS_out2", fill=TEAL_BG, edge=TEAL, radius=0.10)


def slide4(s) -> None:
    """Text left; the running dashboard on the right, as evidence not decoration."""
    clear_ours(s)
    move(s, "TextBox 8", 0.50, 1.26, 7.05, 5.35)
    x, w = 7.78, 5.15
    live = crop(RAW_LIVE, CROP_LIVE, (18, 88, 1345, 478))
    surf = crop(RAW_MAP, CROP_MAP, (30, 252, 1336, 902))
    textbox(s, x, 1.28, w, 0.24,
            [("THE SYSTEM AS IT RUNS TODAY", 10, True, BLUE)],
            name="AS_shdr", align=PP_ALIGN.LEFT)
    with Image.open(live) as im:
        h = w * im.height / im.width
    s.shapes.add_picture(str(live), Inches(x), Inches(1.56), Inches(w), Inches(h)) \
        .name = "AS_shot1"
    y2 = 1.56 + h + 0.13
    with Image.open(surf) as im:
        h2 = w * im.height / im.width
    s.shapes.add_picture(str(surf), Inches(x), Inches(y2), Inches(w), Inches(h2)) \
        .name = "AS_shot2"
    y3 = y2 + h2 + 0.07
    textbox(s, x, y3, w, 0.32,
            [("Live 72 h forecast with intervals and GRAP probability; "
              "wind-aware downscaled surface over the 77-station network. "
              "Both read the local cache.", 8, False, MUTED)],
            name="AS_scap", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    stats = [("9.1 M", "rows cached"), ("77", "NCR stations"),
             ("740 d", "CAMS archive"), ("₹0", "data cost")]
    sw = (w - 3 * 0.10) / 4
    for i, (big, small) in enumerate(stats):
        node(s, x + i * (sw + 0.10), y3 + 0.40, sw, 0.52,
             [(big, 12, True, INK), (small, 8, False, MUTED)],
             name=f"AS_stat{i}", fill=PALE, edge=EDGE, radius=0.14)


def slide5(s) -> None:
    """Text left; the result the whole project rests on, drawn to scale."""
    clear_ours(s)
    move(s, "TextBox 8", 0.50, 1.26, 7.35, 5.35)
    x, w = 8.10, 4.85
    node(s, x, 1.40, w, 0.40,
         [("72 h forecast error, 5 rolling folds", 11, True, WHITE)],
         name="AS_chdr", fill=INK, edge=INK, radius=0.16)
    textbox(s, x, 1.86, w, 0.22,
            [("RMSE in µg/m³ — lower is better", 9, False, MUTED)],
            name="AS_csub")

    bars = [("Raw CAMS", 91.4, RGBColor(0xB0, 0xBE, 0xCC), INK),
            ("Persistence", 77.8, AMBER, INK),
            ("Airshed", 61.9, BLUE, INK)]
    lx, bx, span = x + 0.02, x + 1.28, 2.72   # label column, bar origin, full scale
    top, bh, gap = 2.20, 0.46, 0.26
    for i, (label, val, col, _) in enumerate(bars):
        y = top + i * (bh + gap)
        textbox(s, lx, y, 1.20, bh, [(label, 10, i == 2, INK if i == 2 else MUTED)],
                name=f"AS_blab{i}", align=PP_ALIGN.RIGHT)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx),
                                 Inches(y + 0.06), Inches(span * val / 91.4),
                                 Inches(bh - 0.12))
        bar.name = f"AS_bar{i}"
        bar.adjustments[0] = 0.30
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        bar.shadow.inherit = False
        bar.text_frame.text = ""
        textbox(s, bx + span * val / 91.4 + 0.06, y, 0.75, bh,
                [(f"{val:.1f}", 10, i == 2, INK)], name=f"AS_bval{i}",
                align=PP_ALIGN.LEFT)

    y = top + 3 * (bh + gap) + 0.06
    node(s, x, y, w, 0.52,
         [("−31.5% vs raw CAMS   ·   −20.6% vs persistence", 11, True, TEAL),
          ("better on 5 of 5 seasonal folds", 8.5, False, TEAL)],
         name="AS_gain", fill=TEAL_BG, edge=TEAL, radius=0.14)
    node(s, x, y + 0.64, w, 0.52,
         [("GRAP Stage III: 73.8% recall at 72 h lead", 11, True, AMBER),
          ("per-class recall, never overall accuracy — Stage IV is rare",
           8.5, False, AMBER)],
         name="AS_grap", fill=AMBER_BG, edge=AMBER, radius=0.14)
    half = (w - 0.12) / 2
    node(s, x, y + 1.28, half, 0.62,
         [("85.1% interval coverage", 9.5, True, INK),
          ("against an 80% target", 8, False, MUTED)],
         name="AS_cov", fill=PALE, edge=EDGE, radius=0.12)
    node(s, x + half + 0.12, y + 1.28, half, 0.62,
         [("81.8 µg/m³ held-out, spatial", 9.5, True, INK),
          ("leave-one-station-out, 24.6% better than IDW", 8, False, MUTED)],
         name="AS_loso", fill=PALE, edge=EDGE, radius=0.12)
    textbox(s, x, y + 2.02, w, 0.62,
            [("Gain vs raw CAMS by fold: +12.4 / +33.7 / +18.7 / +35.2 / +57.5%",
              8.5, False, MUTED),
             ("Held-out episode, 22 Dec 2025 — observed 259 µg/m³, raw CAMS off "
              "by 123, Airshed off by 91.", 8.5, False, MUTED)],
            name="AS_note", anchor=MSO_ANCHOR.TOP)


def slide6(s) -> None:
    """References left; what is actually on disk, and how to re-check it, right."""
    clear_ours(s)
    move(s, "TextBox 8", 0.50, 1.32, 7.85, 5.20)
    x, w = 8.62, 4.30
    node(s, x, 1.42, w, 0.40,
         [("EVIDENCE ON DISK", 11, True, WHITE)],
         name="AS_ehdr", fill=INK, edge=INK, radius=0.16)
    rows = [("CPCB ground truth", "2,552 d · 1.54 M rows"),
            ("CAMS archive", "740 d · 1.37 M rows"),
            ("GFS meteorology archive", "710 d · 1.21 M rows"),
            ("Meteorology, lead-matched", "553 d · 2.70 M rows"),
            ("METAR visibility, VIDP", "740 d · 18.9 k obs"),
            ("FIRMS fire detections", "274 d · 23.5 k rows")]
    rh = 0.36
    for i, (a, b) in enumerate(rows):
        y = 1.94 + i * (rh + 0.06)
        node(s, x, y, w, rh,
             [(f"{a}   ·   {b}", 9, False, INK)],
             name=f"AS_row{i}", fill=PALE, edge=EDGE, radius=0.16)
    y = 1.94 + len(rows) * (rh + 0.06) + 0.06
    node(s, x, y, w, 0.82,
         [("Verified, not asserted", 10, True, TEAL),
          ("airshed gate rebuilds a week of every feature with sockets "
           "physically blocked, so \"reads from cache\" is a test, not a claim.",
           8.5, False, TEAL)],
         name="AS_gate", fill=TEAL_BG, edge=TEAL, radius=0.10)
    node(s, x, y + 0.92, w, 0.82,
         [("Next, and already scheduled", 10, True, AMBER),
          ("November 2026 — a second episode season settles fires, the upwind "
           "corridor, coupling and Stage IV in one run.", 8.5, False, AMBER)],
         name="AS_next", fill=AMBER_BG, edge=AMBER, radius=0.10)


def layout(stage1: Path) -> None:
    prs = Presentation(str(stage1))
    for fn, idx in ((slide2, 1), (slide3, 2), (slide4, 3), (slide5, 4), (slide6, 5)):
        fn(prs.slides[idx])
    DECK.unlink(missing_ok=True)
    prs.save(str(DECK))
    print(f"  wrote {DECK.relative_to(ROOT)} ({DECK.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    stage1 = fill_text()
    layout(stage1)
    shutil.rmtree(WORK, ignore_errors=True)
    stage1.unlink(missing_ok=True)
